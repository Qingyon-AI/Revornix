import asyncio
import base64
import inspect
import json
import re
import time
import uuid
from datetime import datetime, timezone
from io import BytesIO

import crud
from langfuse import propagate_attributes
from langfuse.openai import AsyncOpenAI

from base_implement.image_generate_engine_base import ImageGenerateEngineBase
from common.ai import (
    _get_user_ai_interaction_language,
    build_structured_output_language_instruction,
)
from common.logger import exception_logger, info_logger
from common.markdown_helpers import get_markdown_content_by_section_id
from common.usage_billing import persist_model_usage_from_completion
from data.sql.base import async_session_context
from prompts.section_ppt import PPT_SCRIPT_SYSTEM, build_ppt_script_user_prompt
from proxy.ai_model_proxy import AIModelProxy
from proxy.engine_proxy import EngineProxy
from proxy.file_system_proxy import FileSystemProxy
from protocol.remote_file_service import RemoteFileServiceProtocol
from schemas.section import PptPlanResult, PptSlidePlan
from workflow.cancelled import WorkflowCancelledError

try:
    from pptx import Presentation
    from pptx.util import Inches
except Exception:
    Presentation = None
    Inches = None

PPT_GENERATE_CONCURRENCY = 3
PPT_SCRIPT_MAX_CHARS = 18_000
PPT_MANIFEST_VERSION = 2
PPT_SLIDE_GENERATE_MAX_ATTEMPTS = 3
PPT_SLIDE_GENERATE_RETRY_DELAY_SECONDS = 1.2


def _get_section_ppt_manifest_path(section_id: int) -> str:
    return f"generated/sections/{section_id}/ppt/manifest.json"


def _normalize_datetime(value: datetime | None = None) -> str:
    return (value or datetime.now(timezone.utc)).isoformat()


async def _ensure_section_ppt_not_cancelled(
    *,
    remote_file_service: RemoteFileServiceProtocol,
    manifest_path: str,
    section_id: int,
) -> None:
    try:
        raw_content = await remote_file_service.get_file_content_by_file_path(manifest_path)
    except Exception:
        return

    manifest = json.loads(raw_content.decode("utf-8") if isinstance(raw_content, bytes) else raw_content)
    if isinstance(manifest, dict) and manifest.get("status") == "cancelled":
        raise WorkflowCancelledError(f"Section PPT task cancelled: section_id={section_id}")


def _compact_markdown(markdown: str, max_chars: int) -> str:
    normalized = markdown.strip()
    if len(normalized) <= max_chars:
        return normalized
    head_limit = max_chars // 2
    tail_limit = max_chars - head_limit - 32
    return (
        normalized[:head_limit].rstrip()
        + "\n\n...\n\n"
        + normalized[-tail_limit:].lstrip()
    )


async def _safe_close_async_client(client: AsyncOpenAI) -> None:
    close_fn = getattr(client, "close", None)
    if callable(close_fn):
        try:
            result = close_fn()
            if inspect.isawaitable(result):
                await result
        except Exception as e:
            exception_logger.warning(f"[SectionPPT] close llm client failed: {e}")
        return

    aclose_fn = getattr(client, "aclose", None)
    if callable(aclose_fn):
        try:
            result = aclose_fn()
            if inspect.isawaitable(result):
                await result
        except Exception as e:
            exception_logger.warning(f"[SectionPPT] aclose llm client failed: {e}")


async def _write_manifest(
    *,
    remote_file_service: RemoteFileServiceProtocol,
    manifest_path: str,
    payload: dict,
) -> None:
    await remote_file_service.upload_raw_content_to_path(
        file_path=manifest_path,
        content=json.dumps(payload, ensure_ascii=False, indent=2).encode("utf-8"),
        content_type="application/json",
    )


def _extract_image_payload(image_markdown: str) -> tuple[str, bytes] | None:
    normalized = image_markdown.strip()
    markdown_match = re.fullmatch(r"!\[[^\]]*\]\((data:image/[^)]+)\)", normalized)
    if markdown_match is not None:
        normalized = markdown_match.group(1)

    data_match = re.fullmatch(
        r"(data:image/(?P<mime_subtype>[^;]+);base64,(?P<payload>[A-Za-z0-9+/=]+))",
        normalized,
    )
    if data_match is None:
        return None

    mime_subtype = data_match.group("mime_subtype").lower()
    extension = "jpg" if mime_subtype == "jpeg" else mime_subtype.replace("+xml", "")
    payload = data_match.group("payload").strip()
    missing_padding = len(payload) % 4
    if missing_padding:
        payload += "=" * (4 - missing_padding)
    try:
        content = base64.b64decode(payload, validate=True)
    except Exception as decode_error:
        exception_logger.warning(
            f"[SectionPPT] failed to decode generated image payload: {decode_error}"
        )
        return None

    return extension, content


async def _persist_slide_image(
    *,
    remote_file_service: RemoteFileServiceProtocol,
    section_id: int,
    slide_id: str,
    image_markdown: str,
) -> tuple[str | None, bytes | None]:
    payload = _extract_image_payload(image_markdown)
    if payload is None:
        return None, None

    extension, content = payload
    content_type = (
        "image/jpeg"
        if extension == "jpg"
        else "image/svg+xml"
        if extension == "svg"
        else f"image/{extension}"
    )
    file_path = f"images/sections/{section_id}/ppt/{slide_id}-{uuid.uuid4().hex}.{extension}"
    await remote_file_service.upload_raw_content_to_path(
        file_path=file_path,
        content=content,
        content_type=content_type,
    )
    return file_path, content


def _build_pptx_bytes(slide_images: list[bytes]) -> bytes | None:
    if Presentation is None or Inches is None:
        return None
    if not slide_images:
        return None

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank_layout = presentation.slide_layouts[6]

    for image_content in slide_images:
        slide = presentation.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            BytesIO(image_content),
            0,
            0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _parse_theme_prompt(data: dict) -> str | None:
    """Convert the new theme object or legacy theme_prompt string to a plain string."""
    theme = data.get("theme")
    if isinstance(theme, dict):
        parts: list[str] = []
        if theme.get("color_palette"):
            parts.append(f"Colors: {theme['color_palette']}")
        if theme.get("visual_style"):
            parts.append(f"Style: {theme['visual_style']}")
        return ". ".join(parts) if parts else None
    # Fall back to legacy field
    return data.get("theme_prompt") or None


def _parse_slide_plan(raw: dict) -> PptSlidePlan:
    """Map a raw slide dict (new script schema or legacy schema) to PptSlidePlan."""
    # image_prompt is the rich prompt written by the script LLM;
    # fall back to legacy "prompt" key if absent.
    image_prompt = (raw.get("image_prompt") or raw.get("prompt") or "").strip()
    # speaker_notes doubles as the human-readable summary
    speaker_notes = raw.get("speaker_notes") or raw.get("summary") or ""
    return PptSlidePlan(
        id=raw.get("id", ""),
        title=raw.get("title", ""),
        summary=speaker_notes,
        prompt=image_prompt,
        slide_type=raw.get("type") or raw.get("slide_type"),
        key_points=raw.get("key_points") or [],
        speaker_notes=speaker_notes or None,
        layout=raw.get("layout"),
    )


async def _plan_section_ppt(
    *,
    user_id: int,
    markdown: str,
    model_id: int | None = None,
) -> PptPlanResult:
    async with async_session_context() as db:
        db_user = await crud.user.get_user_by_id_async(db=db, user_id=user_id)
    if db_user is None:
        raise RuntimeError("User not found")
    resolved_model_id = model_id or db_user.default_document_reader_model_id
    if resolved_model_id is None:
        raise RuntimeError("Default document reader model not set")

    script_markdown = _compact_markdown(markdown, PPT_SCRIPT_MAX_CHARS)
    user_prompt = build_ppt_script_user_prompt(markdown=script_markdown)
    language_instruction = build_structured_output_language_instruction(
        await _get_user_ai_interaction_language(user_id),
    )
    model_conf = (
        await AIModelProxy.create(
            user_id=user_id,
            model_id=resolved_model_id,
        )
    ).get_configuration()

    with propagate_attributes(
        user_id=str(user_id),
        tags=[f"model:{model_conf.model_name}"],
    ):
        client = AsyncOpenAI(
            api_key=model_conf.api_key,
            base_url=model_conf.base_url,
        )
        try:
            completion = await client.chat.completions.create(
                model=model_conf.model_name,
                messages=[
                    {
                        "role": "system",
                        "content": f"{PPT_SCRIPT_SYSTEM}\n\n{language_instruction}",
                    },
                    {"role": "user", "content": user_prompt},
                ],
                response_format={"type": "json_object"},
            )
            await persist_model_usage_from_completion(
                user_id=user_id,
                model_id=resolved_model_id,
                completion=completion,
                source="plan_section_ppt",
            )
            resp_text = completion.choices[0].message.content
            if not resp_text:
                raise RuntimeError("PPT script generation returned empty response")
            data = json.loads(resp_text)
            if not isinstance(data, dict):
                raise RuntimeError("PPT script response is not a JSON object")

            raw_slides = data.get("slides") or []
            if not isinstance(raw_slides, list):
                raw_slides = []
            slides = [_parse_slide_plan(s) for s in raw_slides[:7]]
            if len(slides) < 4:
                raise RuntimeError("PPT script returned too few slides")

            return PptPlanResult(
                title=data.get("title"),
                subtitle=data.get("subtitle"),
                theme_prompt=_parse_theme_prompt(data),
                slides=slides,
            )
        finally:
            await _safe_close_async_client(client)


async def _generate_slide_image(
    *,
    index: int,
    engine: ImageGenerateEngineBase,
    slide: PptSlidePlan,
    semaphore: asyncio.Semaphore,
) -> tuple[int, PptSlidePlan, str | None]:
    async with semaphore:
        last_error: Exception | None = None
        for attempt in range(1, PPT_SLIDE_GENERATE_MAX_ATTEMPTS + 1):
            started_at = time.perf_counter()
            try:
                image = await engine.generate_image(slide.prompt)
            except Exception as e:
                last_error = e
                elapsed_ms = (time.perf_counter() - started_at) * 1000
                exception_logger.warning(
                    f"[SectionPPT] slide generate failed: slide_id={slide.id}, "
                    f"attempt={attempt}, elapsed_ms={elapsed_ms:.2f}, error={e}"
                )
                if attempt < PPT_SLIDE_GENERATE_MAX_ATTEMPTS:
                    await asyncio.sleep(
                        PPT_SLIDE_GENERATE_RETRY_DELAY_SECONDS * attempt
                    )
                    continue
                raise

            elapsed_ms = (time.perf_counter() - started_at) * 1000
            if image:
                info_logger.info(
                    f"[SectionPPT] slide generated: slide_id={slide.id}, "
                    f"attempt={attempt}, elapsed_ms={elapsed_ms:.2f}"
                )
                return index, slide, image

            exception_logger.warning(
                f"[SectionPPT] slide generate returned empty payload: slide_id={slide.id}, "
                f"attempt={attempt}, elapsed_ms={elapsed_ms:.2f}"
            )
            if attempt < PPT_SLIDE_GENERATE_MAX_ATTEMPTS:
                await asyncio.sleep(PPT_SLIDE_GENERATE_RETRY_DELAY_SECONDS * attempt)

        if last_error is not None:
            raise last_error
        return index, slide, None


async def run_section_ppt_workflow(
    *,
    section_id: int,
    user_id: int,
    model_id: int | None = None,
    image_engine_id: int | None = None,
) -> None:
    remote_file_service = await FileSystemProxy.create(user_id=user_id)
    manifest_path = _get_section_ppt_manifest_path(section_id)
    now = _normalize_datetime()
    existing_manifest: dict | None = None

    async with async_session_context() as db:
        db_user = await crud.user.get_user_by_id_async(db=db, user_id=user_id)
        db_section = await crud.section.get_section_by_section_id_async(db=db, section_id=section_id)
        if db_user is None:
            raise RuntimeError("User not found")
        if db_section is None:
            raise RuntimeError("Section not found")
        resolved_image_engine_id = image_engine_id or db_user.default_image_generate_engine_id
        if resolved_image_engine_id is None:
            raise RuntimeError("Default image generate engine not set")

    try:
        raw_manifest = await remote_file_service.get_file_content_by_file_path(manifest_path)
        parsed_manifest = json.loads(
            raw_manifest.decode("utf-8") if isinstance(raw_manifest, bytes) else raw_manifest
        )
        if isinstance(parsed_manifest, dict):
            existing_manifest = parsed_manifest
    except Exception:
        existing_manifest = None

    initial_manifest = {
        "version": PPT_MANIFEST_VERSION,
        "status": "processing",
        "celery_task_id": (
            existing_manifest.get("celery_task_id")
            if isinstance(existing_manifest, dict)
            else None
        ),
        "title": None,
        "subtitle": None,
        "theme_prompt": None,
        "pptx_file_name": None,
        "error_message": None,
        "create_time": now,
        "update_time": now,
        "slides": [],
    }
    await _write_manifest(
        remote_file_service=remote_file_service,
        manifest_path=manifest_path,
        payload=initial_manifest,
    )
    await _ensure_section_ppt_not_cancelled(
        remote_file_service=remote_file_service,
        manifest_path=manifest_path,
        section_id=section_id,
    )

    try:
        markdown = await get_markdown_content_by_section_id(
            section_id=section_id,
            user_id=user_id,
            remote_file_service=remote_file_service,
            allow_missing=False,
        )
        plan = await _plan_section_ppt(
            user_id=user_id,
            markdown=markdown,
            model_id=model_id,
        )
        await _ensure_section_ppt_not_cancelled(
            remote_file_service=remote_file_service,
            manifest_path=manifest_path,
            section_id=section_id,
        )
        manifest = {
            **initial_manifest,
            "title": plan.title,
            "subtitle": plan.subtitle,
            "theme_prompt": plan.theme_prompt,
            "slides": [
                {
                    "id": slide.id,
                    "type": slide.slide_type,
                    "title": slide.title,
                    "key_points": slide.key_points,
                    "speaker_notes": slide.speaker_notes,
                    "layout": slide.layout,
                    "summary": slide.summary,
                    "prompt": slide.prompt,
                    "image_file_name": None,
                }
                for slide in plan.slides
            ],
            "update_time": _normalize_datetime(),
        }
        await _write_manifest(
            remote_file_service=remote_file_service,
            manifest_path=manifest_path,
            payload=manifest,
        )

        engine = await EngineProxy.create_image_generate_engine(
            user_id=user_id,
            engine_id=resolved_image_engine_id,
        )
        semaphore = asyncio.Semaphore(
            min(PPT_GENERATE_CONCURRENCY, max(1, len(plan.slides)))
        )
        slide_image_payloads: list[bytes] = []

        slide_results: list[tuple[int, bytes]] = []
        tasks = [
            _generate_slide_image(
                index=index,
                engine=engine,
                slide=slide,
                semaphore=semaphore,
            )
            for index, slide in enumerate(plan.slides)
        ]
        for result in asyncio.as_completed(tasks):
            index, slide, image_markdown = await result
            await _ensure_section_ppt_not_cancelled(
                remote_file_service=remote_file_service,
                manifest_path=manifest_path,
                section_id=section_id,
            )
            if not image_markdown:
                raise RuntimeError(
                    f"Slide image generation returned empty response: {slide.id}"
                )

            image_file_name, image_content = await _persist_slide_image(
                remote_file_service=remote_file_service,
                section_id=section_id,
                slide_id=slide.id,
                image_markdown=image_markdown,
            )
            if image_file_name is None or image_content is None:
                raise RuntimeError(f"Slide image payload is invalid: {slide.id}")

            slide_results.append((index, image_content))
            manifest["slides"][index]["image_file_name"] = image_file_name
            manifest["update_time"] = _normalize_datetime()
            await _write_manifest(
                remote_file_service=remote_file_service,
                manifest_path=manifest_path,
                payload=manifest,
            )

        slide_results.sort(key=lambda item: item[0])
        slide_image_payloads = [content for _, content in slide_results]

        pptx_content = _build_pptx_bytes(slide_image_payloads)
        if pptx_content is not None:
            pptx_file_name = f"generated/sections/{section_id}/ppt/{uuid.uuid4().hex}.pptx"
            await remote_file_service.upload_raw_content_to_path(
                file_path=pptx_file_name,
                content=pptx_content,
                content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
            manifest["pptx_file_name"] = pptx_file_name

        manifest["status"] = "success"
        manifest["update_time"] = _normalize_datetime()
        await _write_manifest(
            remote_file_service=remote_file_service,
            manifest_path=manifest_path,
            payload=manifest,
        )
    except WorkflowCancelledError:
        cancelled_manifest = {
            **initial_manifest,
            "status": "cancelled",
            "celery_task_id": None,
            "error_message": "Cancelled by user",
            "update_time": _normalize_datetime(),
        }
        await _write_manifest(
            remote_file_service=remote_file_service,
            manifest_path=manifest_path,
            payload=cancelled_manifest,
        )
        raise
    except Exception as e:
        exception_logger.error(f"[SectionPPT] workflow failed: section={section_id}, error={e}")
        failure_manifest = {
            **initial_manifest,
            "status": "failed",
            "celery_task_id": None,
            "error_message": str(e),
            "update_time": _normalize_datetime(),
        }
        try:
            await _write_manifest(
                remote_file_service=remote_file_service,
                manifest_path=manifest_path,
                payload=failure_manifest,
            )
        except Exception as manifest_error:
            exception_logger.error(
                f"[SectionPPT] failed to write failure manifest: section={section_id}, error={manifest_error}"
            )
